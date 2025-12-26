"""
File I/O service for reading and writing files safely.

Handles:
- Encoding detection
- Line ending normalization
- Atomic writes
- Permission handling
"""

from __future__ import annotations

import os
import shutil
import tempfile
import logging
from dataclasses import dataclass
from enum import Enum, auto
from pathlib import Path
from typing import BinaryIO, Iterator, Optional

import chardet
import pypdf # Added for PDF processing
import docx # Added for Word processing
import openpyxl # Added for Excel processing
import pptx # Added for PowerPoint processing


class LineEnding(Enum):
    """Line ending style."""
    LF = auto()      # Unix: \n
    CRLF = auto()    # Windows: \r\n
    CR = auto()      # Old Mac: \r
    MIXED = auto()   # Mixed endings
    NONE = auto()    # No line endings (single line or binary)


@dataclass
class FileContent:
    """Container for file content with metadata."""
    content: str
    lines: list[str]
    encoding: str
    line_ending: LineEnding
    bom: bool
    size: int
    
    @property
    def line_count(self) -> int:
        return len(self.lines)


@dataclass
class ReadResult:
    """Result of a file read operation."""
    success: bool
    content: Optional[FileContent] = None
    error: Optional[str] = None
    is_binary: bool = False


@dataclass
class WriteResult:
    """Result of a file write operation."""
    success: bool
    bytes_written: int = 0
    error: Optional[str] = None


class FileIOService:
    """Service for safe file I/O operations."""
    
    # Common text encodings to try
    ENCODINGS = ['utf-8', 'utf-8-sig', 'utf-16', 'utf-16-le', 'utf-16-be', 
                 'ascii', 'iso-8859-1', 'cp1252', 'latin-1']
    
    # Binary file signatures (magic bytes)
    BINARY_SIGNATURES = [
        b'\x00',           # Null byte (strong indicator)
        b'\x89PNG',        # PNG
        b'\xff\xd8\xff',   # JPEG
        b'GIF8',           # GIF
        b'PK\x03\x04',     # ZIP
        b'\x1f\x8b',       # GZIP
        b'%PDF',           # PDF
        b'\x7fELF',        # ELF
        b'MZ',             # Windows executable
    ]
    
    def __init__(
        self,
        default_encoding: str = 'utf-8',
        fallback_encoding: str = 'latin-1',
        binary_check_size: int = 8192
    ):
        self.default_encoding = default_encoding
        self.fallback_encoding = fallback_encoding
        self.binary_check_size = binary_check_size
    
    def read_file(
        self,
        path: Path | str,
        encoding: Optional[str] = None,
        normalize_line_endings: bool = False,
        max_text_size: int = 50 * 1024 * 1024  # 50MB default limit
    ) -> ReadResult:
        """
        Read a text file with automatic encoding detection.
        
        Args:
            path: Path to the file
            encoding: Force specific encoding (auto-detect if None)
            normalize_line_endings: Convert all line endings to \n
            max_text_size: Maximum file size in bytes to read as text
            
        Returns:
            ReadResult with content or error information
        """
        path = Path(path)
        
        if not path.exists():
            return ReadResult(success=False, error=f"File not found: {path}")
        
        if not path.is_file():
            return ReadResult(success=False, error=f"Not a file: {path}")

        # Check file size before reading
        try:
            file_size = path.stat().st_size
            if file_size > max_text_size:
                return ReadResult(
                    success=False, 
                    error=f"File too large for text comparison ({file_size / 1024 / 1024:.2f} MB). Max size is {max_text_size / 1024 / 1024:.2f} MB."
                )
        except OSError as e:
            return ReadResult(success=False, error=f"Could not check file size: {e}")
        
        # Handle PDF files specifically
        if path.suffix.lower() == '.pdf':
            pdf_text = self._extract_text_from_pdf(path)
            if pdf_text is not None:
                # Treat extracted PDF text as a regular text file for comparison
                lines = self._split_lines_preserve_endings(pdf_text)
                return ReadResult(
                    success=True,
                    content=FileContent(
                        content=pdf_text,
                        lines=lines,
                        encoding='utf-8', # Assuming extracted text is always UTF-8
                        line_ending=self._detect_line_ending(pdf_text),
                        bom=False,
                        size=len(pdf_text.encode('utf-8'))
                    )
                )
            else:
                return ReadResult(success=False, error=f"Failed to extract text from PDF: {path}")
        
        # Handle Word (.docx) files specifically
        if path.suffix.lower() == '.docx':
            docx_text = self._extract_text_from_docx(path)
            if docx_text is not None:
                # Treat extracted DOCX text as a regular text file for comparison
                lines = self._split_lines_preserve_endings(docx_text)
                return ReadResult(
                    success=True,
                    content=FileContent(
                        content=docx_text,
                        lines=lines,
                        encoding='utf-8', # Assuming extracted text is always UTF-8
                        line_ending=self._detect_line_ending(docx_text),
                        bom=False,
                        size=len(docx_text.encode('utf-8'))
                    )
                )
            else:
                return ReadResult(success=False, error=f"Failed to extract text from DOCX: {path}")

        # Handle Excel (.xlsx) files specifically
        if path.suffix.lower() == '.xlsx':
            xlsx_text = self._extract_text_from_xlsx(path)
            if xlsx_text is not None:
                # Treat extracted XLSX text as a regular text file for comparison
                lines = self._split_lines_preserve_endings(xlsx_text)
                return ReadResult(
                    success=True,
                    content=FileContent(
                        content=xlsx_text,
                        lines=lines,
                        encoding='utf-8', # Assuming extracted text is always UTF-8
                        line_ending=self._detect_line_ending(xlsx_text),
                        bom=False,
                        size=len(xlsx_text.encode('utf-8'))
                    )
                )
            else:
                return ReadResult(success=False, error=f"Failed to extract text from XLSX: {path}")

        # Handle PowerPoint (.pptx) files specifically
        if path.suffix.lower() == '.pptx':
            pptx_text = self._extract_text_from_pptx(path)
            if pptx_text is not None:
                # Treat extracted PPTX text as a regular text file for comparison
                lines = self._split_lines_preserve_endings(pptx_text)
                return ReadResult(
                    success=True,
                    content=FileContent(
                        content=pptx_text,
                        lines=lines,
                        encoding='utf-8', # Assuming extracted text is always UTF-8
                        line_ending=self._detect_line_ending(pptx_text),
                        bom=False,
                        size=len(pptx_text.encode('utf-8'))
                    )
                )
            else:
                return ReadResult(success=False, error=f"Failed to extract text from PPTX: {path}")



        
        try:
            # Check if binary
            if self._is_binary_file(path):
                return ReadResult(success=False, is_binary=True, 
                                error="File appears to be binary")
            
            # Read raw bytes
            raw_content = path.read_bytes()
            size = len(raw_content)
            
            # Detect encoding
            detected_encoding = encoding or self._detect_encoding(raw_content)
            
            # Check for BOM
            bom = False
            if raw_content.startswith(b'\xef\xbb\xbf'):
                bom = True
                detected_encoding = 'utf-8-sig'
            elif raw_content.startswith(b'\xff\xfe'):
                bom = True
                detected_encoding = 'utf-16-le'
            elif raw_content.startswith(b'\xfe\xff'):
                bom = True
                detected_encoding = 'utf-16-be'
            
            # Decode content
            try:
                content = raw_content.decode(detected_encoding)
            except (UnicodeDecodeError, LookupError):
                # Fallback
                content = raw_content.decode(self.fallback_encoding, errors='replace')
                detected_encoding = self.fallback_encoding
            
            # Detect line endings
            line_ending = self._detect_line_ending(content)
            
            # Split into lines
            if normalize_line_endings:
                content = content.replace('\r\n', '\n').replace('\r', '\n')
                lines = content.split('\n')
            else:
                lines = self._split_lines_preserve_endings(content)
            
            return ReadResult(
                success=True,
                content=FileContent(
                    content=content,
                    lines=lines,
                    encoding=detected_encoding,
                    line_ending=line_ending,
                    bom=bom,
                    size=size
                )
            )
            
        except PermissionError:
            return ReadResult(success=False, error=f"Permission denied: {path}")
        except OSError as e:
            return ReadResult(success=False, error=f"OS error: {e}")
    
    def read_file_lines(
        self,
        path: Path | str,
        encoding: Optional[str] = None
    ) -> Iterator[str]:
        """
        Read a file line by line (memory efficient for large files).
        
        Yields lines with line endings stripped.
        """
        path = Path(path)
        encoding = encoding or self._detect_encoding_quick(path)
        
        try:
            with open(path, 'r', encoding=encoding, errors='replace') as f:
                for line in f:
                    yield line
        except Exception as e:
            logging.error(f"FileIOService - Failed to read lines from {path}: {e}")
            return
    
    def write_file(
        self,
        path: Path | str,
        content: str | list[str],
        encoding: str = 'utf-8',
        line_ending: LineEnding = LineEnding.LF,
        atomic: bool = True,
        create_backup: bool = False
    ) -> WriteResult:
        """
        Write content to a file.
        
        Args:
            path: Path to write to
            content: String content or list of lines
            encoding: Encoding to use
            line_ending: Line ending style
            atomic: Use atomic write (write to temp then move)
            create_backup: Create .bak backup of existing file
            
        Returns:
            WriteResult with success status
        """
        path = Path(path)
        
        # Convert lines to string
        if isinstance(content, list):
            line_sep = self._get_line_separator(line_ending)
            content = line_sep.join(content)
            if not content.endswith(line_sep):
                content += line_sep
        
        try:
            # Create backup if requested
            if create_backup and path.exists():
                backup_path = path.with_suffix(path.suffix + '.bak')
                shutil.copy2(path, backup_path)
            
            # Encode content
            encoded = content.encode(encoding)
            
            if atomic:
                # Write to temporary file then move
                dir_path = path.parent
                dir_path.mkdir(parents=True, exist_ok=True)
                
                fd, temp_path = tempfile.mkstemp(dir=dir_path)
                try:
                    os.write(fd, encoded)
                    os.close(fd)
                    shutil.move(temp_path, path)
                except Exception:
                    os.close(fd)
                    if os.path.exists(temp_path):
                        os.unlink(temp_path)
                    raise
            else:
                path.parent.mkdir(parents=True, exist_ok=True)
                path.write_bytes(encoded)
            
            return WriteResult(success=True, bytes_written=len(encoded))
            
        except PermissionError:
            return WriteResult(success=False, error=f"Permission denied: {path}")
        except OSError as e:
            return WriteResult(success=False, error=f"OS error: {e}")
    
    def read_binary(
        self,
        path: Path | str,
        chunk_size: int = 65536
    ) -> Iterator[bytes]:
        """
        Read binary file in chunks.
        
        Yields chunks of bytes for memory-efficient processing.
        """
        path = Path(path)
        
        with open(path, 'rb') as f:
            while chunk := f.read(chunk_size):
                yield chunk
    
    def compare_files_binary(
        self,
        path1: Path | str,
        path2: Path | str,
        chunk_size: int = 65536
    ) -> tuple[bool, Optional[int]]:
        """
        Compare two files byte-by-byte.
        
        Returns:
            Tuple of (are_identical, first_difference_offset)
        """
        path1, path2 = Path(path1), Path(path2)
        
        # Quick size check
        if path1.stat().st_size != path2.stat().st_size:
            return False, 0
        
        offset = 0
        with open(path1, 'rb') as f1, open(path2, 'rb') as f2:
            while True:
                chunk1 = f1.read(chunk_size)
                chunk2 = f2.read(chunk_size)
                
                if chunk1 != chunk2:
                    # Find exact offset
                    for i, (b1, b2) in enumerate(zip(chunk1, chunk2)):
                        if b1 != b2:
                            return False, offset + i
                    # Different lengths
                    return False, offset + min(len(chunk1), len(chunk2))
                
                if not chunk1:  # EOF
                    break
                    
                offset += len(chunk1)
        
        return True, None
    
    def _extract_text_from_pdf(self, path: Path) -> Optional[str]:
        """Extract text content from a PDF file."""
        try:
            reader = pypdf.PdfReader(path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            logging.error(f"Error extracting text from PDF {path}: {e}")
            return None
            
    def _extract_text_from_docx(self, path: Path) -> Optional[str]:
        """Extract text content from a Word (.docx) file."""
        try:
            document = docx.Document(path)
            text = []
            for paragraph in document.paragraphs:
                text.append(paragraph.text)
            return "\n".join(text)
        except Exception as e:
            logging.error(f"Error extracting text from DOCX {path}: {e}")
            return None
            
    def _extract_text_from_xlsx(self, path: Path) -> Optional[str]:
        """Extract text content from an Excel (.xlsx) file."""
        try:
            workbook = openpyxl.load_workbook(path, data_only=True)
            full_text = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                full_text.append(f"===== Sheet: {sheet_name} =====")
                for row in sheet.iter_rows():
                    row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                    full_text.append("\t".join(row_values)) # Use tab as delimiter for columns
                full_text.append("\n") # Add a blank line after each sheet for readability
            return "\n".join(full_text)
        except Exception as e:
            logging.error(f"Error extracting text from XLSX {path}: {e}")
            return None
            
    def _extract_text_from_pptx(self, path: Path) -> Optional[str]:
        """Extract text content from a PowerPoint (.pptx) file."""
        try:
            presentation = pptx.Presentation(path)
            full_text = []
            for i, slide in enumerate(presentation.slides):
                full_text.append(f"===== Slide {i+1} =====")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
                full_text.append("\n") # Add a blank line after each slide for readability
            return "\n".join(full_text)
        except Exception as e:
            logging.error(f"Error extracting text from PPTX {path}: {e}")
            return None
    
    def _is_binary_file(self, path: Path) -> bool:
        """Check if a file is binary."""
        if not path or not path.exists() or not path.is_file():
            return False
        try:
            with open(path, 'rb') as f:
                chunk = f.read(self.binary_check_size)
            
            # Check for binary signatures
            for sig in self.BINARY_SIGNATURES:
                if chunk.startswith(sig):
                    return True
            
            # Check for null bytes
            if b'\x00' in chunk:
                return True
            
            # Check ratio of non-text bytes
            non_text = sum(1 for b in chunk if b < 9 or (b > 13 and b < 32))
            if len(chunk) > 0 and non_text / len(chunk) > 0.3:
                return True
            
            return False
            
        except Exception:
            return True  # Assume binary on error
    
    def _detect_encoding(self, content: bytes) -> str:
        """Detect encoding of content."""
        if not content:
            return self.default_encoding
        
        # Use chardet for detection
        result = chardet.detect(content)
        
        if result['confidence'] > 0.7 and result['encoding']:
            encoding = result['encoding'].lower()
            # Normalize encoding names
            if encoding == 'ascii':
                return 'utf-8'  # ASCII is subset of UTF-8
            return encoding
        
        return self.default_encoding
    
    def _detect_encoding_quick(self, path: Path) -> str:
        """Quick encoding detection from file header."""
        try:
            with open(path, 'rb') as f:
                header = f.read(4096)
            return self._detect_encoding(header)
        except Exception:
            return self.default_encoding
    
    def _detect_line_ending(self, content: str) -> LineEnding:
        """Detect line ending style in content."""
        crlf_count = content.count('\r\n')
        lf_count = content.count('\n') - crlf_count
        cr_count = content.count('\r') - crlf_count
        
        if crlf_count == 0 and lf_count == 0 and cr_count == 0:
            return LineEnding.NONE
        
        total = crlf_count + lf_count + cr_count
        
        if crlf_count == total:
            return LineEnding.CRLF
        elif lf_count == total:
            return LineEnding.LF
        elif cr_count == total:
            return LineEnding.CR
        else:
            return LineEnding.MIXED
    
    def _split_lines_preserve_endings(self, content: str) -> list[str]:
        """Split content into lines, preserving line endings."""
        lines = []
        current = []
        
        i = 0
        while i < len(content):
            char = content[i]
            current.append(char)
            
            if char == '\n':
                lines.append(''.join(current))
                current = []
            elif char == '\r':
                if i + 1 < len(content) and content[i + 1] == '\n':
                    current.append('\n')
                    i += 1
                lines.append(''.join(current))
                current = []
            
            i += 1
        
        if current:
            lines.append(''.join(current))
        
        return lines
    
    def _get_line_separator(self, line_ending: LineEnding) -> str:
        """Get the line separator string for a line ending type."""
        if line_ending == LineEnding.CRLF:
            return '\r\n'
        elif line_ending == LineEnding.CR:
            return '\r'
        else:
            return '\n'


class TempFileManager:
    """Manager for temporary files with automatic cleanup."""
    
    def __init__(self, prefix: str = "compare_", suffix: str = ""):
        self.prefix = prefix
        self.suffix = suffix
        self._temp_files: list[Path] = []
        self._temp_dirs: list[Path] = []
    
    def create_temp_file(
        self,
        content: str = "",
        encoding: str = 'utf-8'
    ) -> Path:
        """Create a temporary file with optional content."""
        fd, path = tempfile.mkstemp(prefix=self.prefix, suffix=self.suffix)
        path = Path(path)
        
        if content:
            os.write(fd, content.encode(encoding))
        os.close(fd)
        
        self._temp_files.append(path)
        return path
    
    def create_temp_dir(self) -> Path:
        """Create a temporary directory."""
        path = Path(tempfile.mkdtemp(prefix=self.prefix))
        self._temp_dirs.append(path)
        return path
    
    def cleanup(self):
        """Clean up all temporary files and directories."""
        for path in self._temp_files:
            try:
                if path.exists():
                    path.unlink()
            except Exception as e:
                logging.warning(f"TempFileManager - Cleanup failed for temp file {path}: {e}")
        
        for path in self._temp_dirs:
            try:
                if path.exists():
                    shutil.rmtree(path)
            except Exception as e:
                logging.warning(f"TempFileManager - Cleanup failed for temp directory {path}: {e}")
        
        self._temp_files.clear()
        self._temp_dirs.clear()
    
    def __enter__(self):
        return self
    
    def __exit__(self, *args):
        self.cleanup()